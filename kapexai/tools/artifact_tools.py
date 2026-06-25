"""Presentation, chart, and diagram tools for artifact-producing agents."""

import re
from pathlib import Path
from typing import Any, Literal

from langchain.tools import tool
from pptx import Presentation

from kapexai.tools.charts import (
    create_bar_chart,
    create_donut_chart,
    create_line_chart,
    create_waterfall_chart,
)
from kapexai.tools.runtime import post_text


OUTPUT_ROOT = Path("output").resolve()


def _safe_output_path(folder: str, filename: str, suffix: str) -> Path:
    safe_name = re.sub(r"[^A-Za-z0-9._-]+", "_", filename).strip("._")
    if not safe_name:
        raise ValueError("filename must contain a letter or number")
    if not safe_name.lower().endswith(suffix):
        safe_name += suffix
    path = (OUTPUT_ROOT / folder / safe_name).resolve()
    if OUTPUT_ROOT not in path.parents:
        raise ValueError("output path must remain inside the output directory")
    path.parent.mkdir(parents=True, exist_ok=True)
    return path


@tool
def create_finance_chart(
    chart_type: Literal["line", "bar", "stacked_bar", "waterfall", "donut"],
    title: str,
    labels: list[str],
    series: dict[str, list[float]],
    filename: str,
    subtitle: str | None = None,
    value_style: Literal["currency", "number", "percent"] = "currency",
) -> dict:
    """Create a polished PNG chart for reports or presentations.

    Waterfall expects one series of changes. Donut expects one series of values.
    """
    output_path = _safe_output_path("charts", filename, ".png")
    if chart_type == "line":
        result = create_line_chart(
            labels,
            series,
            output_path,
            title=title,
            subtitle=subtitle,
            value_style=value_style,
        )
    elif chart_type in {"bar", "stacked_bar"}:
        result = create_bar_chart(
            labels,
            series,
            output_path,
            title=title,
            subtitle=subtitle,
            value_style=value_style,
            stacked=chart_type == "stacked_bar",
        )
    elif chart_type == "waterfall":
        if len(series) != 1:
            raise ValueError("waterfall requires exactly one change series")
        result = create_waterfall_chart(
            labels,
            next(iter(series.values())),
            output_path,
            title=title,
            subtitle=subtitle,
            value_style=value_style,
        )
    else:
        if len(series) != 1:
            raise ValueError("donut requires exactly one value series")
        result = create_donut_chart(
            labels,
            next(iter(series.values())),
            output_path,
            title=title,
            subtitle=subtitle,
        )
    return {
        "chart_type": chart_type,
        "output_path": str(Path(result).resolve()),
        "title": title,
    }


@tool
def create_presentation_from_outline(
    title: str,
    subtitle: str,
    slides: list[dict[str, Any]],
    filename: str,
) -> dict:
    """Build a simple PowerPoint from an approved slide outline.

    Each slide supports title and bullets. This tool does not invent content.
    """
    output_path = _safe_output_path("presentations", filename, ".pptx")
    presentation = Presentation()
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title_slide.shapes.title.text = title
    title_slide.placeholders[1].text = subtitle

    for slide_spec in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = str(slide_spec.get("title", ""))
        body = slide.shapes.placeholders[1].text_frame
        body.clear()
        bullets = slide_spec.get("bullets", [])
        for index, bullet in enumerate(bullets):
            paragraph = (
                body.paragraphs[0] if index == 0 else body.add_paragraph()
            )
            if isinstance(bullet, dict):
                paragraph.text = str(bullet.get("text", ""))
                paragraph.level = int(bullet.get("level", 0))
            else:
                paragraph.text = str(bullet)

    presentation.save(output_path)
    return {
        "output_path": str(output_path),
        "slide_count": len(presentation.slides),
        "title": title,
    }


@tool
def validate_mermaid(diagram_source: str) -> dict:
    """Perform fast structural checks on Mermaid source before rendering."""
    stripped = diagram_source.strip()
    allowed_starts = (
        "flowchart ",
        "graph ",
        "sequenceDiagram",
        "stateDiagram",
        "stateDiagram-v2",
        "classDiagram",
        "erDiagram",
        "journey",
        "gantt",
        "timeline",
        "mindmap",
        "quadrantChart",
    )
    issues = []
    if not stripped.startswith(allowed_starts):
        issues.append("Unsupported or missing Mermaid diagram declaration")
    if len(stripped) > 50000:
        issues.append("Diagram exceeds the 50,000 character safety limit")
    if "<script" in stripped.lower() or "javascript:" in stripped.lower():
        issues.append("Unsafe script content is not allowed")
    bracket_pairs = {"[": "]", "(": ")", "{": "}"}
    for opening, closing in bracket_pairs.items():
        if stripped.count(opening) != stripped.count(closing):
            issues.append(f"Unbalanced {opening}{closing} delimiters")
    return {
        "valid": not issues,
        "issues": issues,
        "line_count": len(stripped.splitlines()),
        "character_count": len(stripped),
    }


@tool
def render_mermaid_svg(diagram_source: str, filename: str) -> dict:
    """Render Mermaid source to SVG using the Kroki public HTTP API."""
    validation = validate_mermaid.invoke({"diagram_source": diagram_source})
    if not validation["valid"]:
        raise ValueError("; ".join(validation["issues"]))
    svg = post_text(
        "https://kroki.io/mermaid/svg",
        diagram_source,
        headers={
            "Accept": "image/svg+xml",
            "Content-Type": "text/plain",
        },
        ttl_seconds=86400,
    )
    output_path = _safe_output_path("diagrams", filename, ".svg")
    output_path.write_text(svg, encoding="utf-8")
    return {
        "output_path": str(output_path),
        "format": "svg",
        "source": "https://kroki.io/mermaid/svg",
        "validation": validation,
    }
